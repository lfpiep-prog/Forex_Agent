
import os
from reportlab.pdfgen import canvas
from docx import Document
from pptx import Presentation
import pandas as pd
import xlsxwriter

OUTPUT_DIR = "test_documents"
if not os.path.exists(OUTPUT_DIR):
    os.makedirs(OUTPUT_DIR)

def create_pdf():
    path = os.path.join(OUTPUT_DIR, "test.pdf")
    c = canvas.Canvas(path)
    c.drawString(100, 750, "PDF Generation Verified!")
    c.save()
    print(f"Created {path}")

def create_docx():
    path = os.path.join(OUTPUT_DIR, "test.docx")
    doc = Document()
    doc.add_heading('Word Generation Verified', 0)
    doc.add_paragraph('This confirms python-docx is working.')
    doc.save(path)
    print(f"Created {path}")

def create_pptx():
    path = os.path.join(OUTPUT_DIR, "test.pptx")
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "PowerPoint Verified"
    subtitle.text = "python-pptx is working!"
    prs.save(path)
    print(f"Created {path}")

def create_excel():
    path = os.path.join(OUTPUT_DIR, "test.xlsx")
    df = pd.DataFrame({'Data': [10, 20, 30, 40], 'Labels': ['A', 'B', 'C', 'D']})
    # Using xlsxwriter engine
    writer = pd.ExcelWriter(path, engine='xlsxwriter')
    df.to_excel(writer, sheet_name='Sheet1', index=False)
    writer.close()
    print(f"Created {path}")

if __name__ == "__main__":
    print("Verifying Document Generation Libraries...")
    try:
        create_pdf()
        create_docx()
        create_pptx()
        create_excel()
        print("\nAll checks passed successfully.")
    except Exception as e:
        print(f"\nFAILED: {e}")
